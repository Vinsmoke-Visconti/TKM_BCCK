import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Tu dong xac dinh thu muc goc cua du an
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)

def add_slide_with_image(prs, title_text, bullet_points, img_path=None):
    """Helper function to create a slide with text on the left and image on the right"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    
    # Text box (left side)
    content = slide.placeholders[1]
    content.text = "\n".join(bullet_points)
    
    # Adjust content box size to make room for image
    content.width = Inches(4.5)
    
    # Add image if exists
    if img_path and os.path.exists(img_path):
        try:
            prs.slides[-1].shapes.add_picture(img_path, Inches(4.8), Inches(1.5), width=Inches(4.8))
        except Exception as e:
            print(f"  [WARN] Khong the add anh {img_path}: {e}")
    else:
        # Placeholder or warning if image missing
        txBox = slide.shapes.add_textbox(Inches(5.5), Inches(3), Inches(3), Inches(1))
        tf = txBox.text_frame
        tf.text = f"[Thieu anh: {os.path.basename(img_path) if img_path else 'N/A'}]"

def create_presentation():
    prs = Presentation()
    # Duong dan thu muc anh
    img_dir = os.path.join(PROJECT_ROOT, "baocao", "image")
    docs_dir = os.path.join(PROJECT_ROOT, "docs")
    save_path = os.path.join(PROJECT_ROOT, "baocao", "thuyet_minh_do_an.pptx")

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "BÁO CÁO ĐỒ ÁN THIẾT KẾ MẠNG"
    subtitle.text = "Đề tài: Thiết kế và triển khai mạng Metro Ethernet sử dụng MPLS\nSinh viên: Huỳnh Nguyễn Quốc Việt - 52300267\nNgành: Mạng máy tính và truyền thông dữ liệu"

    # Slide 2: Muc tieu do an
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Mục tiêu và Phạm vi nghiên cứu"
    content = slide.placeholders[1]
    content.text = (
        "- Xây dựng mô hình Metro Ethernet MAN kết nối đa chi nhánh.\n"
        "- Triển khai kỹ thuật chuyển mạch nhãn MPLS (Multiprotocol Label Switching).\n"
        "- So sánh hiệu năng giữa 3 kiến trúc LAN: Flat, 3-Tier và Leaf-Spine.\n"
        "- Đánh giá khả năng mở rộng và tính bảo mật của mạng lõi ISP."
    )

    # Slide 3: So do Topology (Minh chung thiet ke)
    topo_img = os.path.join(docs_dir, "network_detailed_topology.png")
    add_slide_with_image(prs, "2. Thiết kế Topology hệ thống", [
        "- ISP Core: Router P1, P2 chạy OSPF/LDP.",
        "- ISP Edge: Router PE1, PE2 kết nối khách hàng.",
        "- Site A: Flat Network (Mô hình doanh nghiệp nhỏ).",
        "- Site B: Core-Dist-Access (Mô hình tiêu chuẩn).",
        "- Site C: Leaf-Spine (Mô hình Data Center)."
    ], topo_img)

    # Slide 4: Khoi tao he thong (Minh chung Mininet)
    mininet_img = os.path.join(img_dir, "mininet_nodes_net.png")
    add_slide_with_image(prs, "3. Hiện thực hóa trên Mininet", [
        "- Sử dụng LinuxRouter và OVS Switches.",
        "- Cấu hình STP cho Site C để chống Loop lớp 2.",
        "- Quy hoạch IP khoa học: /30 cho Backbone, /24 cho LAN.",
        "- Kiểm tra Nodes/Net: Đầy đủ 11 Switch, 7 Router và 10 Hosts."
    ], mininet_img)

    # Slide 5: Minh chung Dinh tuyen (OSPF & MPLS)
    ospf_img = os.path.join(img_dir, "ospf_neighbor.png")
    add_slide_with_image(prs, "4. Minh chứng Hội tụ định tuyến", [
        "- OSPF đạt trạng thái FULL giữa P và PE.",
        "- LDP phát nhãn tự động cho các dải Loopback và LAN.",
        "- Bảng định tuyến xuất hiện 'Label' (24, 25, implicit-null).",
        "- Chứng minh dữ liệu đi qua hầm MPLS thay vì IP thuần."
    ], ospf_img)

    # Slide 6: Minh chung Thong mach & Bao mat
    trace_img = os.path.join(img_dir, "traceroute_mpls.png")
    add_slide_with_image(prs, "5. Kiểm tra kết nối & Tính ẩn danh", [
        "- Ping xuyên chi nhánh: Đạt 0% Loss, RTT ổn định.",
        "- Traceroute: Xuất hiện dấu * * * tại các hop mạng lõi.",
        "- Giải thích: Tính năng ẩn IP của MPLS giúp bảo mật hạ tầng ISP."
    ], trace_img)

    # Slide 7: Phan tich hieu nang (Bieu do Throughput)
    throughput_img = os.path.join(img_dir, "chart_throughput.png")
    add_slide_with_image(prs, "6. Đánh giá Thông lượng (Throughput)", [
        "- Site B -> Site C (Nội bộ PE): Đạt băng thông cao nhất (~189Mbps).",
        "- Site A -> Site B (Xuyên lõi): Bị giới hạn bởi nút thắt WAN (~100Mbps).",
        "- Leaf-Spine thể hiện sự ổn định vượt trội khi có lưu lượng cao."
    ], throughput_img)

    # Slide 8: Phan tich hieu nang (Bieu do Delay)
    delay_img = os.path.join(img_dir, "chart_delay.png")
    add_slide_with_image(prs, "7. Phân tích Độ trễ (RTT)", [
        "- Độ trễ trung bình duy trì mức 20ms - 27ms.",
        "- MPLS giúp giảm Overhead xử lý tại Core Router.",
        "- Kiến trúc LAN phân tầng giúp tăng tính sẵn sàng (Availability)."
    ], delay_img)

    # Slide 9: Ket luan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "8. Tổng kết đồ án"
    content = slide.placeholders[1]
    content.text = (
        "- Hoàn thành 100% mục tiêu thiết kế và triển khai.\n"
        "- Hệ thống tự động hóa hoàn toàn từ cấu hình đến đo lường.\n"
        "- Nắm vững công nghệ chuyển mạch nhãn hiện đại nhất hiện nay.\n"
        "- Hướng phát triển: Triển khai QoS và Traffic Engineering trên MPLS."
    )

    prs.save(save_path)
    print(f"Presentation updated with evidence images: {save_path}")

if __name__ == "__main__":
    create_presentation()
